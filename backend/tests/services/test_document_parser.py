import io
from unittest.mock import patch

from app.core.upload_types import resolve_document_file_type
from app.services.document_conversion import prepare_uploaded_document
from app.services.document_parser import (
    _XLS_OLE2_MAGIC,
    DocumentParserService,
    _clean_pdf_pages,
    _iter_sheet_rows,
    _normalize_xls_cell,
)


def test_clean_pdf_pages_removes_repeated_margin_boilerplate():
    pages = [
        (
            "DocuSign Envelope ID: ABC\n"
            "Shared report title\n"
            "Section heading\n"
            "Important body line A\n"
            "Detail A1\n"
            "Detail A2\n"
            "Detail A3\n"
            "Detail A4\n"
            "Detail A5\n"
            "6",
            1,
        ),
        (
            "DocuSign Envelope ID: ABC\n"
            "Shared report title\n"
            "Section heading\n"
            "Important body line B\n"
            "Detail B1\n"
            "Detail B2\n"
            "Detail B3\n"
            "Detail B4\n"
            "Detail B5\n"
            "7",
            2,
        ),
        (
            "DocuSign Envelope ID: ABC\n"
            "Shared report title\n"
            "Section heading\n"
            "Important body line C\n"
            "Detail C1\n"
            "Detail C2\n"
            "Detail C3\n"
            "Detail C4\n"
            "Detail C5\n"
            "8",
            3,
        ),
    ]

    cleaned = _clean_pdf_pages(pages)

    assert cleaned == [
        ("Section heading\nImportant body line A\nDetail A1\nDetail A2\nDetail A3\nDetail A4\nDetail A5", 1),
        ("Section heading\nImportant body line B\nDetail B1\nDetail B2\nDetail B3\nDetail B4\nDetail B5", 2),
        ("Section heading\nImportant body line C\nDetail C1\nDetail C2\nDetail C3\nDetail C4\nDetail C5", 3),
    ]


def test_clean_pdf_pages_preserves_repeated_body_labels():
    pages = [
        (
            "Repeated header\n"
            "Expected outputs\n"
            "Target: 10 companies\n"
            "Repeated footer",
            1,
        ),
        (
            "Repeated header\n"
            "Expected outputs\n"
            "Target: 20 companies\n"
            "Repeated footer",
            2,
        ),
        (
            "Repeated header\n"
            "Expected outputs\n"
            "Target: 30 companies\n"
            "Repeated footer",
            3,
        ),
    ]

    cleaned = _clean_pdf_pages(pages)

    assert cleaned == [
        ("Expected outputs\nTarget: 10 companies", 1),
        ("Expected outputs\nTarget: 20 companies", 2),
        ("Expected outputs\nTarget: 30 companies", 3),
    ]


def test_resolve_document_file_type_accepts_pptx_and_iwork_extensions():
    assert resolve_document_file_type("", "deck.pptx") == "pptx"
    assert resolve_document_file_type("", "brief.pages") == "pages"
    assert resolve_document_file_type("", "slides.keynote") == "keynote"
    assert resolve_document_file_type("", "model.numbers") == "numbers"
    assert resolve_document_file_type("", "model.dwg") is None


def test_prepare_uploaded_document_converts_pages_to_docx_filename():
    with patch(
        "app.services.document_conversion._convert_with_libreoffice",
        return_value=b"converted",
    ):
        prepared = prepare_uploaded_document(b"raw", "Brief.pages", "pages")

    assert prepared.content == b"converted"
    assert prepared.filename == "Brief.docx"
    assert prepared.file_type == "docx"


def test_prepare_uploaded_document_converts_keynote_to_pptx_filename():
    with patch(
        "app.services.document_conversion._convert_with_libreoffice",
        return_value=b"converted",
    ):
        prepared = prepare_uploaded_document(b"raw", "Deck.key", "keynote")

    assert prepared.content == b"converted"
    assert prepared.filename == "Deck.pptx"
    assert prepared.file_type == "pptx"


def test_prepare_uploaded_document_converts_numbers_to_xlsx_filename():
    with patch(
        "app.services.document_conversion._convert_with_libreoffice",
        return_value=b"converted",
    ):
        prepared = prepare_uploaded_document(b"raw", "Budget.numbers", "numbers")

    assert prepared.content == b"converted"
    assert prepared.filename == "Budget.xlsx"
    assert prepared.file_type == "xlsx"


def test_parse_xlsx_reads_modern_workbook():
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet.append(["Year", "Rating"])
    sheet.append([2015, "Low"])

    buf = io.BytesIO()
    workbook.save(buf)

    text = DocumentParserService().parse_xlsx(buf.getvalue())

    assert "[Sheet: Summary]" in text
    assert "Year\tRating" in text
    assert "2015\tLow" in text


def test_iter_sheet_rows_routes_legacy_xls_to_xlrd(monkeypatch):
    """Legacy BIFF .xls must not reach openpyxl, which rejects it as "not a zip file"."""

    class _Sheet:
        name = "Risk Tool v4.0 - Template"
        nrows = 2

        def row_values(self, idx):
            return [["Year", "Rating"], [2015.0, 3.5]][idx]

    class _Book:
        def sheets(self):
            return [_Sheet()]

    import xlrd

    monkeypatch.setattr(xlrd, "open_workbook", lambda **_kwargs: _Book())

    sheets = _iter_sheet_rows(_XLS_OLE2_MAGIC + b"biff-payload")

    assert sheets == [("Risk Tool v4.0 - Template", [("Year", "Rating"), (2015, 3.5)])]


def test_normalize_xls_cell_keeps_whole_numbers_readable():
    # xlrd types every number as float; "2015.0" in extracted text is noise.
    assert _normalize_xls_cell(2015.0) == 2015
    assert _normalize_xls_cell(3.5) == 3.5
    assert _normalize_xls_cell("Low") == "Low"
    assert _normalize_xls_cell(None) is None


def test_parse_pptx_extracts_slide_text():
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Project Overview"
    textbox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    textbox.text = "Revenue case\nBase case"

    buf = io.BytesIO()
    presentation.save(buf)

    text = DocumentParserService().parse_pptx(buf.getvalue())

    assert "[Slide 1]" in text
    assert "Project Overview" in text
    assert "Revenue case" in text
