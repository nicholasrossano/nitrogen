import io
import math
import re
import tiktoken

from app.config import get_settings

settings = get_settings()


class DocumentParserService:
    """Service for parsing documents and chunking text"""

    def __init__(self):
        self.chunk_size = settings.chunk_size
        self.chunk_overlap = settings.chunk_overlap
        self.tokenizer = tiktoken.get_encoding("cl100k_base")

    # ------------------------------------------------------------------
    # Plain-text extraction (used for embeddings / backward compat)
    # ------------------------------------------------------------------

    def parse_pdf(self, content: bytes) -> str:
        """Parse PDF content to plain text (all pages concatenated)."""
        return "\n\n".join(text for text, _ in self.parse_pdf_pages(content))

    def parse_pdf_pages(self, content: bytes) -> list[tuple[str, int]]:
        """Parse PDF returning (text, page_number) pairs (1-indexed)."""
        import pdfplumber

        pages: list[tuple[str, int]] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    pages.append((page_text, page.page_number))
        return _clean_pdf_pages(pages)

    def parse_docx(self, content: bytes) -> str:
        """Parse DOCX content to plain text."""
        from docx import Document

        doc = Document(io.BytesIO(content))
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        return "\n\n".join(text_parts)

    def parse_xlsx(self, content: bytes) -> str:
        """Parse XLSX/XLS spreadsheet content to plain text."""
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        sheet_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                sheet_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        return "\n\n".join(sheet_parts)

    def parse_pptx(self, content: bytes) -> str:
        """Parse PPTX presentation content to plain text."""
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        slide_parts: list[str] = []

        for idx, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    text_parts.append(text.strip())

                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            rows.append("\t".join(cells))
                    if rows:
                        text_parts.append("\n".join(rows))

            if text_parts:
                slide_parts.append(f"[Slide {idx}]\n" + "\n\n".join(text_parts))

        return "\n\n".join(slide_parts)

    # ------------------------------------------------------------------
    # Rich (HTML) extraction – for content_html column
    # ------------------------------------------------------------------

    def parse_docx_html(self, content: bytes) -> str:
        """Convert DOCX to HTML preserving bold, italic, links, tables, lists."""
        import mammoth

        result = mammoth.convert_to_html(io.BytesIO(content))
        return result.value

    def parse_xlsx_html(self, content: bytes) -> str:
        """Convert XLSX to HTML tables (one per sheet)."""
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        html_parts: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_rows = list(ws.iter_rows(values_only=True))
            non_empty = [r for r in all_rows if any(
                (str(c).strip() if c is not None else "") for c in r
            )]
            if not non_empty:
                continue

            table_html = f'<h3>{_esc(sheet_name)}</h3>\n<table>\n'
            for idx, row in enumerate(non_empty):
                tag = "th" if idx == 0 else "td"
                cells = "".join(
                    f"<{tag}>{_esc(str(c) if c is not None else '')}</{tag}>"
                    for c in row
                )
                table_html += f"<tr>{cells}</tr>\n"
            table_html += "</table>"
            html_parts.append(table_html)

        return "\n".join(html_parts)

    def parse_pptx_html(self, content: bytes) -> str:
        """Convert PPTX slide text to simple HTML sections."""
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        html_parts: list[str] = []

        for idx, slide in enumerate(presentation.slides, start=1):
            slide_fragments = [f"<h3>Slide {idx}</h3>"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    paragraphs = [
                        f"<p>{_esc(line.strip())}</p>"
                        for line in text.splitlines()
                        if line.strip()
                    ]
                    slide_fragments.extend(paragraphs)

                if getattr(shape, "has_table", False):
                    table_html = "<table>\n"
                    for row_idx, row in enumerate(shape.table.rows):
                        tag = "th" if row_idx == 0 else "td"
                        cells = "".join(
                            f"<{tag}>{_esc(cell.text.strip())}</{tag}>"
                            for cell in row.cells
                        )
                        table_html += f"<tr>{cells}</tr>\n"
                    table_html += "</table>"
                    slide_fragments.append(table_html)

            if len(slide_fragments) > 1:
                html_parts.append("\n".join(slide_fragments))

        return "\n".join(html_parts)

    # ------------------------------------------------------------------
    # Chunking – plain text (original)
    # ------------------------------------------------------------------

    def chunk_text(self, text: str) -> list[str]:
        """Split text into chunks of approximately chunk_size tokens."""
        tokens = self.tokenizer.encode(text)
        chunks = []
        start = 0

        while start < len(tokens):
            end = start + self.chunk_size
            chunk_tokens = tokens[start:end]
            chunk_text = self.tokenizer.decode(chunk_tokens)
            chunk_text = self._clean_chunk_end(chunk_text)
            if start > 0:
                chunk_text = self._clean_chunk_start(chunk_text)

            if chunk_text.strip():
                chunks.append(chunk_text.strip())

            start = end - self.chunk_overlap

        return chunks

    def chunk_pdf_pages(
        self, pages: list[tuple[str, int]]
    ) -> list[tuple[str, int]]:
        """Chunk PDF text while tracking the originating page number.

        Returns list of (chunk_text, page_number) tuples.  Each chunk is
        tagged with the page it started on.
        """
        chunks: list[tuple[str, int]] = []

        for page_text, page_num in pages:
            page_chunks = self.chunk_text(page_text)
            for c in page_chunks:
                chunks.append((c, page_num))

        return chunks

    # ------------------------------------------------------------------
    # Chunking – HTML-aware (for content_html)
    # ------------------------------------------------------------------

    def chunk_html(self, html: str) -> list[tuple[str, str]]:
        """Chunk HTML into (plain_text, html_fragment) pairs.

        Splits on block-level element boundaries so tags stay balanced.
        Falls back to the plain-text chunker if the HTML has no block tags.
        """

        # Note: "tr" is intentionally excluded so entire <table> blocks stay
        # together as a single fragment rather than being split row-by-row.
        block_tags = {
            "p", "h1", "h2", "h3", "h4", "h5", "h6",
            "table", "ul", "ol", "li", "blockquote", "div", "section",
        }

        fragments = _split_html_blocks(html, block_tags)
        if not fragments:
            plain = _strip_tags(html)
            text_chunks = self.chunk_text(plain)
            return [(c, f"<p>{_esc(c)}</p>") for c in text_chunks]

        results: list[tuple[str, str]] = []
        current_html_parts: list[str] = []
        current_plain_parts: list[str] = []
        current_tokens = 0

        for frag_html in fragments:
            frag_plain = _strip_tags(frag_html).strip()
            frag_tokens = self.count_tokens(frag_plain) if frag_plain else 0

            if frag_tokens > self.chunk_size:
                if current_html_parts:
                    results.append((
                        "\n\n".join(current_plain_parts),
                        "\n".join(current_html_parts),
                    ))
                    current_html_parts = []
                    current_plain_parts = []
                    current_tokens = 0
                sub_chunks = self.chunk_text(frag_plain)
                for sc in sub_chunks:
                    results.append((sc, f"<p>{_esc(sc)}</p>"))
                continue

            if current_tokens + frag_tokens > self.chunk_size and current_html_parts:
                results.append((
                    "\n\n".join(current_plain_parts),
                    "\n".join(current_html_parts),
                ))
                overlap_count = max(1, len(current_html_parts) // 4)
                current_html_parts = current_html_parts[-overlap_count:]
                current_plain_parts = current_plain_parts[-overlap_count:]
                current_tokens = sum(
                    self.count_tokens(p) for p in current_plain_parts
                )

            current_html_parts.append(frag_html)
            if frag_plain:
                current_plain_parts.append(frag_plain)
            current_tokens += frag_tokens

        if current_html_parts:
            results.append((
                "\n\n".join(current_plain_parts),
                "\n".join(current_html_parts),
            ))

        return results

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _clean_chunk_end(self, text: str) -> str:
        """Snap the chunk end back to the nearest sentence boundary."""
        sentence_ends = {'.', '!', '?', '\n'}

        if len(text) > 50:
            for i in range(len(text) - 1, max(len(text) - 100, 0), -1):
                if text[i] in sentence_ends:
                    return text[:i + 1]

        return text

    def _clean_chunk_start(self, text: str) -> str:
        """Snap the chunk start forward to the nearest sentence boundary."""
        sentence_ends = {'.', '!', '?', '\n'}
        search_limit = min(100, len(text) // 2)

        for i in range(search_limit):
            if text[i] in sentence_ends:
                rest = text[i + 1:].lstrip()
                if rest:
                    return rest

        return text

    def count_tokens(self, text: str) -> int:
        """Count tokens in text."""
        return len(self.tokenizer.encode(text))


# ======================================================================
# Assessment-level helpers
# ======================================================================

_TAG_RE = re.compile(r"<[^>]+>")
_DOCUSIGN_RE = re.compile(r"\bDocuSign\s+Envelope\s+ID\b", re.IGNORECASE)
_STANDALONE_PAGE_RE = re.compile(r"^(?:page\s*)?-?\s*\d{1,4}\s*-?$", re.IGNORECASE)


def _strip_tags(html: str) -> str:
    """Remove all HTML tags, returning plain text."""
    return _TAG_RE.sub("", html)


def _esc(text: str) -> str:
    """Minimal HTML escaping."""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def _clean_pdf_pages(pages: list[tuple[str, int]]) -> list[tuple[str, int]]:
    """Remove conservative PDF boilerplate from extracted page text.

    Only repeated margin lines, DocuSign IDs, and standalone page numbers are
    removed. Repeated body text is preserved so tables and recurring labels are
    not accidentally stripped.
    """

    if not pages:
        return []

    repeated_margin_lines = _find_repeated_margin_lines([text for text, _ in pages])
    cleaned_pages: list[tuple[str, int]] = []
    for text, page_number in pages:
        cleaned = _clean_pdf_page_text(text, repeated_margin_lines)
        if cleaned.strip():
            cleaned_pages.append((cleaned, page_number))
    return cleaned_pages


def _find_repeated_margin_lines(page_texts: list[str]) -> set[str]:
    if len(page_texts) < 3:
        return set()

    line_pages: dict[str, set[int]] = {}
    for page_idx, text in enumerate(page_texts):
        lines = [line for line in text.splitlines() if line.strip()]
        margin_band_size = _pdf_margin_band_size(len(lines))
        margin_lines = lines[:margin_band_size] + lines[-margin_band_size:]
        for line in margin_lines:
            normalized = _normalize_pdf_boilerplate_line(line)
            if normalized:
                line_pages.setdefault(normalized, set()).add(page_idx)

    threshold = max(2, math.ceil(len(page_texts) * 0.5))
    return {
        normalized
        for normalized, seen_pages in line_pages.items()
        if len(seen_pages) >= threshold
    }


def _clean_pdf_page_text(text: str, repeated_margin_lines: set[str]) -> str:
    lines = text.splitlines()
    total = len(lines)
    cleaned: list[str] = []

    for idx, line in enumerate(lines):
        stripped = line.strip()
        normalized = _normalize_pdf_boilerplate_line(stripped)
        margin_band_size = _pdf_margin_band_size(total)
        in_margin = idx < margin_band_size or idx >= max(total - margin_band_size, 0)

        if _DOCUSIGN_RE.search(stripped):
            continue
        if in_margin and normalized in repeated_margin_lines:
            continue
        if in_margin and _STANDALONE_PAGE_RE.match(stripped):
            continue

        cleaned.append(line.rstrip())

    return "\n".join(cleaned).strip()


def _normalize_pdf_boilerplate_line(line: str) -> str:
    return re.sub(r"\s+", " ", line).strip().lower()


def _pdf_margin_band_size(line_count: int) -> int:
    if line_count <= 0:
        return 0
    return min(4, max(1, math.ceil(line_count * 0.15)))


def _split_html_blocks(html: str, block_tags: set[str]) -> list[str]:
    """Split HTML string on block-level element boundaries.

    Returns a list of HTML fragments, each roughly corresponding to one
    block element (or a run of inline content between blocks).
    """
    pattern = r"(<(?:" + "|".join(block_tags) + r")[\s>])"
    parts = re.split(pattern, html, flags=re.IGNORECASE)

    fragments: list[str] = []
    buf = ""
    for part in parts:
        stripped = part.strip()
        if not stripped:
            continue
        match = re.match(
            r"<(" + "|".join(block_tags) + r")[\s>]", stripped, re.IGNORECASE
        )
        if match and buf.strip():
            fragments.append(buf.strip())
            buf = part
        else:
            buf += part
    if buf.strip():
        fragments.append(buf.strip())

    return fragments
